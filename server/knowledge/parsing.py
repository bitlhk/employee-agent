"""Structure-aware document parsing for the EA knowledge service."""

from __future__ import annotations

import csv
from dataclasses import dataclass
import json
import os
from pathlib import Path
import re
import zipfile


MAX_DOCUMENT_CHARS = int(os.environ.get("KNOWLEDGE_MAX_DOCUMENT_CHARS", "2000000"))
MAX_PDF_PAGES = int(os.environ.get("KNOWLEDGE_MAX_PDF_PAGES", "2000"))
MAX_SHEET_ROWS = int(os.environ.get("KNOWLEDGE_MAX_SHEET_ROWS", "20000"))
MAX_ARCHIVE_ENTRIES = int(os.environ.get("KNOWLEDGE_MAX_ARCHIVE_ENTRIES", "10000"))
MAX_EXPANDED_BYTES = int(os.environ.get("KNOWLEDGE_MAX_EXPANDED_BYTES", str(200 * 1024 * 1024)))
MAX_ARCHIVE_RATIO = float(os.environ.get("KNOWLEDGE_MAX_ARCHIVE_RATIO", "120"))

_HEADING_RE = re.compile(
    r"^(?:第[一二三四五六七八九十百千万0-9]+[编章节条款]|"
    r"[一二三四五六七八九十]+[、.]|"
    r"\d+(?:\.\d+){0,4}[、.\s]|"
    r"[（(][一二三四五六七八九十0-9]+[）)]|"
    r"【[^】]{1,60}】).{0,80}$"
)


@dataclass(frozen=True)
class ParsedSegment:
    position: str
    text: str
    heading_path: tuple[str, ...] = ()
    page: int | None = None
    content_type: str = "text"


def clean_text(value: str) -> str:
    value = value.replace("\x00", " ")
    value = re.sub(r"[\u0001-\u0008\u000b\u000c\u000e-\u001f]", " ", value)
    value = re.sub(r"[ \t]+", " ", value)
    value = re.sub(r"\n{4,}", "\n\n\n", value)
    return value.strip()


def _bounded_text(value: str, label: str) -> str:
    value = clean_text(value)
    if len(value) > MAX_DOCUMENT_CHARS:
        raise ValueError(f"{label} extracted text exceeds limit")
    return value


def _position(headings: list[str], fallback: str) -> str:
    return " / ".join(headings) if headings else fallback


def _structured_text_segments(text: str, markdown: bool, fallback: str = "正文") -> list[ParsedSegment]:
    headings: list[str] = []
    buffer: list[str] = []
    segments: list[ParsedSegment] = []

    def flush() -> None:
        body = clean_text("\n".join(buffer))
        buffer.clear()
        if body:
            segments.append(ParsedSegment(
                position=_position(headings, fallback),
                heading_path=tuple(headings),
                text=body,
            ))

    for raw_line in text.splitlines():
        line = raw_line.strip()
        markdown_match = re.match(r"^(#{1,6})\s+(.+?)\s*#*$", line) if markdown else None
        plain_heading = not markdown and bool(line) and len(line) <= 90 and bool(_HEADING_RE.match(line))
        if markdown_match or plain_heading:
            flush()
            if markdown_match:
                level = len(markdown_match.group(1))
                title = clean_text(markdown_match.group(2))[:160]
            else:
                level = 1
                title = clean_text(line)[:160]
            headings[:] = headings[: max(0, level - 1)]
            headings.append(title)
            continue
        buffer.append(raw_line)
    flush()
    return segments or [ParsedSegment(position=fallback, text=clean_text(text))]


def _guard_office_archive(source: Path) -> None:
    try:
        with zipfile.ZipFile(source) as archive:
            entries = archive.infolist()
            if len(entries) > MAX_ARCHIVE_ENTRIES:
                raise ValueError("office archive contains too many entries")
            expanded = sum(max(0, entry.file_size) for entry in entries)
            compressed = sum(max(1, entry.compress_size) for entry in entries)
            if expanded > MAX_EXPANDED_BYTES:
                raise ValueError("office archive expands beyond safety limit")
            if expanded / max(1, compressed) > MAX_ARCHIVE_RATIO:
                raise ValueError("office archive compression ratio is unsafe")
            if any(entry.filename.startswith(("/", "\\")) or ".." in Path(entry.filename).parts for entry in entries):
                raise ValueError("office archive contains an unsafe path")
    except zipfile.BadZipFile as exc:
        raise ValueError("invalid office document archive") from exc


def _read_csv(source: Path) -> list[ParsedSegment]:
    text = source.read_text("utf-8", errors="replace")
    rows = list(csv.reader(text.splitlines()))
    if not rows:
        return []
    if len(rows) > MAX_SHEET_ROWS:
        raise ValueError("CSV row count exceeds safety limit")
    header = rows[0]
    segments: list[ParsedSegment] = []
    batch_size = 40
    for start in range(1, len(rows), batch_size):
        batch = rows[start:start + batch_size]
        lines = [" | ".join(header), *[" | ".join(row) for row in batch]]
        end = start + len(batch)
        segments.append(ParsedSegment(
            position=f"第 {start + 1}-{end + 1} 行",
            text=clean_text("\n".join(lines)),
            content_type="table",
        ))
    if len(rows) == 1:
        segments.append(ParsedSegment(position="第 1 行", text=" | ".join(header), content_type="table"))
    return segments


def _read_pdf(source: Path) -> list[ParsedSegment]:
    from pypdf import PdfReader

    reader = PdfReader(str(source))
    if reader.is_encrypted:
        try:
            if reader.decrypt("") == 0:
                raise ValueError("encrypted PDF is not supported")
        except Exception as exc:
            raise ValueError("encrypted PDF is not supported") from exc
    if len(reader.pages) > MAX_PDF_PAGES:
        raise ValueError("PDF page count exceeds safety limit")
    segments: list[ParsedSegment] = []
    total = 0
    for index, page in enumerate(reader.pages, start=1):
        text = clean_text(page.extract_text() or "")
        total += len(text)
        if total > MAX_DOCUMENT_CHARS:
            raise ValueError("PDF extracted text exceeds safety limit")
        if text:
            segments.append(ParsedSegment(position=f"第 {index} 页", text=text, page=index))
    return segments


def _iter_docx_blocks(document):
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield Table(child, document)


def _read_docx(source: Path) -> list[ParsedSegment]:
    from docx import Document
    from docx.table import Table

    _guard_office_archive(source)
    document = Document(str(source))
    headings: list[str] = []
    paragraphs: list[str] = []
    segments: list[ParsedSegment] = []
    table_index = 0
    total = 0

    def add_segment(text: str, content_type: str = "text", suffix: str = "") -> None:
        nonlocal total
        body = clean_text(text)
        if not body:
            return
        total += len(body)
        if total > MAX_DOCUMENT_CHARS:
            raise ValueError("DOCX extracted text exceeds safety limit")
        base = _position(headings, "正文")
        segments.append(ParsedSegment(
            position=f"{base}{suffix}",
            heading_path=tuple(headings),
            text=body,
            content_type=content_type,
        ))

    def flush_paragraphs() -> None:
        add_segment("\n".join(paragraphs))
        paragraphs.clear()

    for block in _iter_docx_blocks(document):
        if isinstance(block, Table):
            flush_paragraphs()
            table_index += 1
            rows = [" | ".join(clean_text(cell.text) for cell in row.cells) for row in block.rows]
            add_segment("\n".join(rows), "table", f" · 表格 {table_index}")
            continue
        text = clean_text(block.text)
        if not text:
            continue
        style_name = str(getattr(block.style, "name", "") or "")
        match = re.search(r"(?:Heading|标题)\s*([1-6])", style_name, re.IGNORECASE)
        if match:
            flush_paragraphs()
            level = int(match.group(1))
            headings[:] = headings[: max(0, level - 1)]
            headings.append(text[:160])
        else:
            paragraphs.append(text)
    flush_paragraphs()
    return segments


def _read_pptx(source: Path) -> list[ParsedSegment]:
    from pptx import Presentation

    _guard_office_archive(source)
    presentation = Presentation(str(source))
    segments: list[ParsedSegment] = []
    total = 0
    for index, slide in enumerate(presentation.slides, start=1):
        texts = [clean_text(shape.text) for shape in slide.shapes if hasattr(shape, "text") and clean_text(shape.text)]
        if not texts:
            continue
        title = texts[0][:160]
        body = clean_text("\n".join(texts))
        total += len(body)
        if total > MAX_DOCUMENT_CHARS:
            raise ValueError("PPTX extracted text exceeds safety limit")
        segments.append(ParsedSegment(
            position=f"第 {index} 页 · {title}",
            heading_path=(title,),
            text=body,
            page=index,
        ))
    return segments


def _read_xlsx(source: Path) -> list[ParsedSegment]:
    from openpyxl import load_workbook

    _guard_office_archive(source)
    workbook = load_workbook(str(source), read_only=True, data_only=True, keep_links=False)
    segments: list[ParsedSegment] = []
    total = 0
    try:
        for worksheet in workbook.worksheets:
            rows: list[list[str]] = []
            for row_index, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
                if row_index > MAX_SHEET_ROWS:
                    raise ValueError("XLSX row count exceeds safety limit")
                values = [clean_text(str(value)) if value is not None else "" for value in row]
                if any(values):
                    rows.append(values)
            if not rows:
                continue
            header = rows[0]
            for start in range(1, len(rows), 40):
                batch = rows[start:start + 40]
                lines = [" | ".join(header), *[" | ".join(row) for row in batch]]
                body = clean_text("\n".join(lines))
                total += len(body)
                if total > MAX_DOCUMENT_CHARS:
                    raise ValueError("XLSX extracted text exceeds safety limit")
                end = start + len(batch)
                segments.append(ParsedSegment(
                    position=f"工作表 {worksheet.title} · 第 {start + 1}-{end + 1} 行",
                    heading_path=(worksheet.title,),
                    text=body,
                    content_type="table",
                ))
            if len(rows) == 1:
                segments.append(ParsedSegment(
                    position=f"工作表 {worksheet.title} · 第 1 行",
                    heading_path=(worksheet.title,),
                    text=" | ".join(header),
                    content_type="table",
                ))
    finally:
        workbook.close()
    return segments


def read_segments(source: Path) -> list[ParsedSegment]:
    ext = source.suffix.lower()
    if ext == ".md":
        text = _bounded_text(source.read_text("utf-8", errors="replace"), "Markdown")
        return [segment for segment in _structured_text_segments(text, markdown=True) if segment.text]
    if ext == ".txt":
        text = _bounded_text(source.read_text("utf-8", errors="replace"), "text document")
        return [segment for segment in _structured_text_segments(text, markdown=False) if segment.text]
    if ext == ".csv":
        return _read_csv(source)
    if ext == ".json":
        raw = source.read_text("utf-8", errors="replace")
        try:
            raw = json.dumps(json.loads(raw), ensure_ascii=False, indent=2)
        except json.JSONDecodeError:
            pass
        return [ParsedSegment(position="正文", text=_bounded_text(raw, "JSON document"))]
    if ext in {".yaml", ".yml"}:
        return [ParsedSegment(position="正文", text=_bounded_text(source.read_text("utf-8", errors="replace"), "YAML document"))]
    if ext == ".pdf":
        return _read_pdf(source)
    if ext == ".docx":
        return _read_docx(source)
    if ext == ".pptx":
        return _read_pptx(source)
    if ext == ".xlsx":
        return _read_xlsx(source)
    raise ValueError(f"unsupported document type: {ext or 'unknown'}")
