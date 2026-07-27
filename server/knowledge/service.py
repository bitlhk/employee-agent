#!/usr/bin/env python3
"""Local-only LlamaIndex retrieval service for EA knowledge bases."""

from __future__ import annotations

import asyncio
from functools import lru_cache
import json
import os
from pathlib import Path
import re
import shutil
import tempfile
from typing import Any, Literal

from fastapi import FastAPI, Header, HTTPException
from dotenv import load_dotenv
from pydantic import BaseModel, Field

from llama_index.core import StorageContext, VectorStoreIndex, load_index_from_storage
from llama_index.core.node_parser import SentenceSplitter
from llama_index.core.schema import Document, NodeWithScore, QueryBundle, TextNode
from llama_index.retrievers.bm25 import BM25Retriever


APP_ROOT = Path(os.environ.get("APP_ROOT") or Path(__file__).resolve().parents[2]).resolve()
load_dotenv(APP_ROOT / ".env", override=False)
DATA_ROOT = (APP_ROOT / "data" / "knowledge").resolve()
INDEX_ROOT = DATA_ROOT / "indexes"
TOKEN_PATH = Path(os.environ.get("KNOWLEDGE_SERVICE_TOKEN_FILE") or DATA_ROOT / ".service-token")
MAX_TOTAL_CHARS = int(os.environ.get("KNOWLEDGE_MAX_TOTAL_CHARS", "8000000"))
MAX_NODES = int(os.environ.get("KNOWLEDGE_MAX_NODES", "12000"))
KB_ID_RE = re.compile(r"^kb_[A-Za-z0-9_-]{8,56}$")

app = FastAPI(title="EA Knowledge Service", docs_url=None, redoc_url=None)
_locks: dict[str, asyncio.Lock] = {}


class SourceDocument(BaseModel):
    id: str = Field(min_length=3, max_length=64)
    name: str = Field(min_length=1, max_length=240)
    path: str = Field(min_length=1, max_length=1200)


class IndexRequest(BaseModel):
    knowledge_base_id: str
    documents: list[SourceDocument] = Field(default_factory=list, max_length=2000)


class SearchRequest(BaseModel):
    knowledge_base_id: str
    query: str = Field(min_length=1, max_length=4000)
    top_k: int = Field(default=6, ge=1, le=20)


class MultiSearchRequest(BaseModel):
    knowledge_base_ids: list[str] = Field(min_length=1, max_length=8)
    query: str = Field(min_length=1, max_length=4000)
    top_k: int = Field(default=4, ge=1, le=12)
    mode: Literal["auto", "forced"] = "auto"


def _token() -> str:
    try:
        return TOKEN_PATH.read_text("utf-8").strip()
    except OSError:
        return ""


def _authorize(value: str | None) -> None:
    expected = _token()
    if not expected or not value or value != expected:
        raise HTTPException(status_code=401, detail="unauthorized")


def _kb_id(value: str) -> str:
    if not KB_ID_RE.fullmatch(value):
        raise HTTPException(status_code=400, detail="invalid knowledge base id")
    return value


def _safe_source_path(value: str) -> Path:
    source = Path(value).resolve(strict=True)
    documents_root = (DATA_ROOT / "documents").resolve()
    if source != documents_root and documents_root not in source.parents:
        raise ValueError("document path is outside knowledge storage")
    if not source.is_file():
        raise ValueError("document is not a file")
    return source


def _clean_text(value: str) -> str:
    value = value.replace("\x00", " ")
    value = re.sub(r"[\u0001-\u0008\u000b\u000c\u000e-\u001f]", " ", value)
    value = re.sub(r"[ \t]+", " ", value)
    value = re.sub(r"\n{4,}", "\n\n\n", value)
    return value.strip()


def _read_segments(source: Path) -> list[tuple[str, str]]:
    ext = source.suffix.lower()
    if ext in {".txt", ".md", ".csv", ".json", ".yaml", ".yml"}:
        return [("正文", _clean_text(source.read_text("utf-8", errors="replace")))]
    if ext == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(str(source))
        return [(f"第 {index + 1} 页", _clean_text(page.extract_text() or "")) for index, page in enumerate(reader.pages)]
    if ext == ".docx":
        from docx import Document

        document = Document(str(source))
        text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        for table in document.tables:
            text += "\n" + "\n".join(" | ".join(cell.text for cell in row.cells) for row in table.rows)
        return [("正文", _clean_text(text))]
    if ext == ".pptx":
        from pptx import Presentation

        presentation = Presentation(str(source))
        segments: list[tuple[str, str]] = []
        for index, slide in enumerate(presentation.slides):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
            segments.append((f"第 {index + 1} 页", _clean_text("\n".join(texts))))
        return segments
    if ext == ".xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(str(source), read_only=True, data_only=True)
        segments = []
        for worksheet in workbook.worksheets:
            lines = []
            for row in worksheet.iter_rows(values_only=True):
                line = " | ".join(str(value) if value is not None else "" for value in row).strip(" |")
                if line:
                    lines.append(line)
            segments.append((f"工作表 {worksheet.title}", _clean_text("\n".join(lines))))
        workbook.close()
        return segments
    raise ValueError(f"unsupported document type: {ext or 'unknown'}")


def _build_nodes(documents: list[SourceDocument]) -> tuple[list[TextNode], dict[str, int]]:
    splitter = SentenceSplitter(chunk_size=700, chunk_overlap=100)
    nodes: list[TextNode] = []
    counts: dict[str, int] = {}
    total_chars = 0
    for document in documents:
        source = _safe_source_path(document.path)
        ordinal = 0
        for position, text in _read_segments(source):
            if not text:
                continue
            total_chars += len(text)
            if total_chars > MAX_TOTAL_CHARS:
                raise ValueError("knowledge base text exceeds indexing limit")
            parent = Document(text=text)
            for chunk in splitter.get_nodes_from_documents([parent], show_progress=False):
                chunk_text = _clean_text(chunk.get_content())
                if not chunk_text:
                    continue
                ordinal += 1
                nodes.append(TextNode(
                    id_=f"{document.id}:{ordinal}",
                    text=chunk_text,
                    metadata={
                        "document_id": document.id,
                        "document_name": document.name,
                        "position": position,
                        "ordinal": ordinal,
                    },
                ))
                if len(nodes) > MAX_NODES:
                    raise ValueError("knowledge base contains too many chunks")
        counts[document.id] = ordinal
    return nodes, counts


@lru_cache(maxsize=1)
def _embedding_model():
    api_key = os.environ.get("KNOWLEDGE_EMBED_API_KEY", "").strip()
    api_base = os.environ.get("KNOWLEDGE_EMBED_API_BASE", "").strip()
    model = os.environ.get("KNOWLEDGE_EMBED_MODEL", "").strip()
    if not api_key or not api_base or not model:
        return None
    from llama_index.embeddings.openai import OpenAIEmbedding

    batch_size = max(1, min(int(os.environ.get("KNOWLEDGE_EMBED_BATCH_SIZE", "32")), 64))
    return OpenAIEmbedding(
        model_name=model,
        api_key=api_key,
        api_base=api_base,
        embed_batch_size=batch_size,
    )


def _persist_nodes(nodes: list[TextNode], target: Path) -> None:
    payload = [node.to_dict() for node in nodes]
    (target / "nodes.json").write_text(json.dumps(payload, ensure_ascii=False), "utf-8")


def _load_nodes(target: Path) -> list[TextNode]:
    payload = json.loads((target / "nodes.json").read_text("utf-8"))
    return [TextNode.from_dict(item) for item in payload]


def _bm25_text(value: str) -> str:
    import jieba

    return " ".join(token.strip().lower() for token in jieba.cut_for_search(value) if token.strip())


def _build_index(request: IndexRequest) -> dict[str, Any]:
    kb_id = _kb_id(request.knowledge_base_id)
    INDEX_ROOT.mkdir(parents=True, exist_ok=True)
    destination = INDEX_ROOT / kb_id
    if not request.documents:
        shutil.rmtree(destination, ignore_errors=True)
        return {"ok": True, "chunk_count": 0, "document_chunks": {}, "vector_enabled": False}

    nodes, counts = _build_nodes(request.documents)
    if not nodes:
        raise ValueError("documents contain no extractable text")

    temp_root = Path(tempfile.mkdtemp(prefix=f".{kb_id}-", dir=INDEX_ROOT))
    try:
        _persist_nodes(nodes, temp_root)
        bm25_dir = temp_root / "bm25"
        bm25_dir.mkdir(parents=True, exist_ok=True)
        bm25_nodes = [TextNode(
            id_=node.node_id,
            text=_bm25_text(node.get_content()),
            metadata=node.metadata,
        ) for node in nodes]
        BM25Retriever.from_defaults(
            nodes=bm25_nodes,
            similarity_top_k=min(20, len(nodes)),
            skip_stemming=True,
            token_pattern=r"(?u)\b\w+\b",
        ).persist(str(bm25_dir))

        vector_enabled = False
        embed_model = _embedding_model()
        if embed_model is not None:
            try:
                import faiss
                from llama_index.vector_stores.faiss import FaissVectorStore

                dimension = len(embed_model.get_text_embedding("EA knowledge dimension probe"))
                vector_store = FaissVectorStore(faiss_index=faiss.IndexFlatL2(dimension))
                storage_context = StorageContext.from_defaults(vector_store=vector_store)
                index = VectorStoreIndex(nodes, storage_context=storage_context, embed_model=embed_model, show_progress=False)
                vector_dir = temp_root / "vector"
                index.storage_context.persist(persist_dir=str(vector_dir))
                (temp_root / "vector.json").write_text(json.dumps({"model": os.environ.get("KNOWLEDGE_EMBED_MODEL", ""), "dimension": dimension}), "utf-8")
                vector_enabled = True
            except Exception as exc:
                (temp_root / "vector-error.txt").write_text(str(exc)[:2000], "utf-8")

        metadata = {
            "knowledge_base_id": kb_id,
            "chunk_count": len(nodes),
            "document_chunks": counts,
            "vector_enabled": vector_enabled,
        }
        (temp_root / "manifest.json").write_text(json.dumps(metadata, ensure_ascii=False, indent=2), "utf-8")
        backup = INDEX_ROOT / f".{kb_id}.old"
        shutil.rmtree(backup, ignore_errors=True)
        if destination.exists():
            destination.rename(backup)
        temp_root.rename(destination)
        shutil.rmtree(backup, ignore_errors=True)
        return {"ok": True, **metadata}
    except Exception:
        shutil.rmtree(temp_root, ignore_errors=True)
        raise


def _rrf(result_sets: list[list[NodeWithScore]], top_k: int) -> list[tuple[TextNode, float]]:
    scores: dict[str, float] = {}
    nodes: dict[str, TextNode] = {}
    for results in result_sets:
        for rank, result in enumerate(results, start=1):
            node = result.node
            node_id = str(node.node_id)
            scores[node_id] = scores.get(node_id, 0.0) + 1.0 / (60 + rank)
            nodes[node_id] = node
    ordered = sorted(scores.items(), key=lambda item: item[1], reverse=True)[:top_k]
    return [(nodes[node_id], score) for node_id, score in ordered]


def _auto_query_candidate(value: str) -> bool:
    compact = re.sub(r"\s+", "", value).lower()
    if not compact or len(compact) < 4:
        return False
    if re.fullmatch(r"(?:你好|您好|嗨|hello|hi|在吗|谢谢|感谢|收到|好的|好|ok|test|测试)[!！。.，,？?]*", compact):
        return False
    if re.fullmatch(r"/(?:new|reset|help|status|tools|model|context|usage|tasks).*", compact):
        return False
    if re.search(r"(?:天气|气温|温度|下雨|降雨|降水|空气质量|台风|几点|现在时间|今天几号|星期几|实时路况)", compact):
        return False
    return True


def _tag_node(node: TextNode, knowledge_base_id: str) -> TextNode:
    node.metadata = {**node.metadata, "knowledge_base_id": knowledge_base_id}
    return node


def _search_indexes(request: MultiSearchRequest) -> dict[str, Any]:
    knowledge_base_ids = list(dict.fromkeys(_kb_id(value) for value in request.knowledge_base_ids))
    targets: list[tuple[str, Path, list[TextNode]]] = []
    for knowledge_base_id in knowledge_base_ids:
        target = INDEX_ROOT / knowledge_base_id
        if not target.is_dir():
            continue
        nodes = _load_nodes(target)
        if nodes:
            targets.append((knowledge_base_id, target, nodes))
    if not targets:
        return {
            "ok": True,
            "triggered": False,
            "results": [],
            "retrieval": "unavailable",
            "metrics": {"knowledge_base_count": 0, "bm25_max_score": 0, "vector_min_distance": None},
        }

    auto_candidate = request.mode == "forced" or _auto_query_candidate(request.query)
    candidate_k = max(request.top_k * 3, request.top_k)
    bm25_candidates: list[NodeWithScore] = []
    bm25_max_score = 0.0
    for knowledge_base_id, target, nodes in targets:
        nodes_by_id = {str(node.node_id): node for node in nodes}
        bm25 = BM25Retriever.from_persist_dir(str(target / "bm25"))
        bm25.similarity_top_k = min(candidate_k, len(nodes))
        for result in bm25.retrieve(QueryBundle(_bm25_text(request.query))):
            score = float(result.score or 0)
            original = nodes_by_id.get(str(result.node.node_id))
            if original is None or score <= 0:
                continue
            bm25_max_score = max(bm25_max_score, score)
            bm25_candidates.append(NodeWithScore(node=_tag_node(original, knowledge_base_id), score=score))
    bm25_candidates.sort(key=lambda item: float(item.score or 0), reverse=True)

    vector_candidates: list[NodeWithScore] = []
    vector_min_distance: float | None = None
    embed_model = _embedding_model()
    if auto_candidate and embed_model is not None and any((target / "vector").is_dir() for _, target, _ in targets):
        try:
            from llama_index.vector_stores.faiss import FaissVectorStore

            query_embedding = embed_model.get_query_embedding(request.query)
            query_bundle = QueryBundle(query_str=request.query, embedding=query_embedding)
            for knowledge_base_id, target, nodes in targets:
                if not (target / "vector").is_dir():
                    continue
                vector_store = FaissVectorStore.from_persist_dir(str(target / "vector"))
                storage_context = StorageContext.from_defaults(persist_dir=str(target / "vector"), vector_store=vector_store)
                index = load_index_from_storage(storage_context, embed_model=embed_model)
                retriever = index.as_retriever(similarity_top_k=min(candidate_k, len(nodes)))
                for result in retriever.retrieve(query_bundle):
                    distance = float(result.score) if result.score is not None else None
                    if distance is not None:
                        vector_min_distance = distance if vector_min_distance is None else min(vector_min_distance, distance)
                    vector_candidates.append(NodeWithScore(node=_tag_node(result.node, knowledge_base_id), score=result.score))
            vector_candidates.sort(key=lambda item: float(item.score if item.score is not None else 1e9))
        except Exception:
            vector_candidates = []
            vector_min_distance = None

    bm25_threshold = float(os.environ.get("KNOWLEDGE_AUTO_BM25_MIN_SCORE", "1.2"))
    vector_threshold = float(os.environ.get("KNOWLEDGE_AUTO_VECTOR_MAX_DISTANCE", "0.88"))
    triggered = request.mode == "forced" or (
        auto_candidate
        and (
            bm25_max_score >= bm25_threshold
            or (vector_min_distance is not None and vector_min_distance <= vector_threshold)
        )
    )
    result_sets: list[list[NodeWithScore]] = []
    if bm25_candidates:
        result_sets.append(bm25_candidates)
    if vector_candidates:
        result_sets.append(vector_candidates)
    if not triggered or not result_sets:
        return {
            "ok": True,
            "triggered": False,
            "results": [],
            "retrieval": "hybrid" if vector_candidates else "bm25",
            "metrics": {
                "knowledge_base_count": len(targets),
                "bm25_max_score": round(bm25_max_score, 6),
                "vector_min_distance": round(vector_min_distance, 6) if vector_min_distance is not None else None,
            },
        }

    results = []
    document_counts: dict[str, int] = {}
    for node, score in _rrf(result_sets, request.top_k * 3):
        knowledge_base_id = str(node.metadata.get("knowledge_base_id") or "")
        document_id = str(node.metadata.get("document_id") or "")
        document_key = f"{knowledge_base_id}:{document_id}"
        if document_counts.get(document_key, 0) >= 2:
            continue
        document_counts[document_key] = document_counts.get(document_key, 0) + 1
        results.append({
            "chunk_id": node.node_id,
            "score": score,
            "text": node.get_content()[:2000],
            "knowledge_base_id": knowledge_base_id,
            "document_id": document_id,
            "document_name": str(node.metadata.get("document_name") or ""),
            "position": str(node.metadata.get("position") or "正文"),
            "ordinal": int(node.metadata.get("ordinal") or 0),
        })
        if len(results) >= request.top_k:
            break
    return {
        "ok": True,
        "triggered": bool(results),
        "results": results,
        "retrieval": "hybrid" if vector_candidates else "bm25",
        "metrics": {
            "knowledge_base_count": len(targets),
            "bm25_max_score": round(bm25_max_score, 6),
            "vector_min_distance": round(vector_min_distance, 6) if vector_min_distance is not None else None,
        },
    }


def _search_index(request: SearchRequest) -> dict[str, Any]:
    payload = _search_indexes(MultiSearchRequest(
        knowledge_base_ids=[request.knowledge_base_id],
        query=request.query,
        top_k=request.top_k,
        mode="forced",
    ))
    if not payload["results"] and not (INDEX_ROOT / _kb_id(request.knowledge_base_id)).is_dir():
        raise FileNotFoundError("knowledge index not found")
    return payload


@app.get("/health")
async def health(x_ea_knowledge_token: str | None = Header(default=None)):
    _authorize(x_ea_knowledge_token)
    return {"ok": True, "engine": "llamaindex", "embedding_configured": _embedding_model() is not None}


@app.post("/index")
async def index(request: IndexRequest, x_ea_knowledge_token: str | None = Header(default=None)):
    _authorize(x_ea_knowledge_token)
    kb_id = _kb_id(request.knowledge_base_id)
    lock = _locks.setdefault(kb_id, asyncio.Lock())
    async with lock:
        try:
            return await asyncio.to_thread(_build_index, request)
        except (OSError, ValueError) as exc:
            raise HTTPException(status_code=422, detail=str(exc)[:1000]) from exc


@app.post("/search")
async def search(request: SearchRequest, x_ea_knowledge_token: str | None = Header(default=None)):
    _authorize(x_ea_knowledge_token)
    try:
        return await asyncio.to_thread(_search_index, request)
    except FileNotFoundError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc
    except (OSError, ValueError) as exc:
        raise HTTPException(status_code=422, detail=str(exc)[:1000]) from exc


@app.post("/search-multi")
async def search_multi(request: MultiSearchRequest, x_ea_knowledge_token: str | None = Header(default=None)):
    _authorize(x_ea_knowledge_token)
    try:
        return await asyncio.to_thread(_search_indexes, request)
    except (OSError, ValueError) as exc:
        raise HTTPException(status_code=422, detail=str(exc)[:1000]) from exc


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="127.0.0.1", port=int(os.environ.get("KNOWLEDGE_SERVICE_PORT", "5191")), access_log=False)
